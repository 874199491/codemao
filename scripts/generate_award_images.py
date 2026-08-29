#!/usr/bin/env python3
"""Generate award PNGs for students with score >= threshold using the award template.

Renders the 优秀学员奖状1.pptx template per student: replaces the name text
box ("学生：") and exports the slide to PNG via PowerPoint COM. Existing PNGs
are skipped (idempotent). Threshold defaults to 80.
"""
import argparse
import json
import os
import sys
import tempfile
from pathlib import Path

try:
    sys.stdout.reconfigure(encoding="utf-8", errors="replace")
except Exception:
    pass


def find_student_textbox(slide):
    """Return the textbox whose text contains 学生, recursively through groups."""
    def walk(container):
        for shape in container:
            if shape.shape_type == 6:  # GROUP
                result = walk(shape.shapes)
                if result is not None:
                    return result
            elif shape.has_text_frame:
                text = "".join(run.text for para in shape.text_frame.paragraphs for run in para.runs)
                if "学生" in text:
                    return shape
        return None
    return walk(slide.shapes)


def render_award(template_path: Path, output_png: Path, student_name: str) -> None:
    from pptx import Presentation
    from pptx.util import Pt

    prs = Presentation(str(template_path))
    slide = prs.slides[0]
    box = find_student_textbox(slide)
    if box is None:
        raise RuntimeError("奖状模板中找不到学生姓名文本框")
    for para in box.text_frame.paragraphs:
        for run in para.runs:
            if "学生" in run.text:
                run.text = run.text.replace("学生：", f"学生：{student_name}")
                run.font.size = Pt(32)

    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp:
        tmp_path = Path(tmp.name)
    try:
        prs.save(str(tmp_path))
        import win32com.client

        app = win32com.client.DispatchEx("PowerPoint.Application")
        try:
            pres = app.Presentations.Open(str(tmp_path), WithWindow=False)
            try:
                pres.Slides(1).Export(str(output_png), "PNG", 1600, 1200)
            finally:
                pres.Close()
        finally:
            app.Quit()
    finally:
        tmp_path.unlink(missing_ok=True)


def main() -> int:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("--source-dir", required=True, type=Path, help="月考反馈助手目录（含模板和已生成奖状）")
    parser.add_argument("--manifest", required=True, type=Path, help="manifest.json（含学生姓名与成绩）")
    parser.add_argument("--threshold", type=float, default=80.0, help="奖状最低展示分（默认 80）")
    parser.add_argument("--limit", type=int, default=0, help="仅生成前 N 个（调试用，0=全部）")
    args = parser.parse_args()

    template = args.source_dir / "优秀学员奖状1.pptx"
    if not template.is_file():
        raise RuntimeError(f"奖状模板不存在：{template}")
    award_dir = args.source_dir / "已生成奖状"
    award_dir.mkdir(parents=True, exist_ok=True)
    manifest = json.loads(args.manifest.read_text(encoding="utf-8"))
    students = [
        s for s in manifest.get("students") or []
        if s.get("score") is not None and float(s["score"]) >= args.threshold
    ]
    students.sort(key=lambda s: (-float(s["score"]), s.get("student_name") or ""))
    if args.limit > 0:
        students = students[: args.limit]

    print(f"奖状阈值：{args.threshold:g} 分；符合学生 {len(students)} 人", flush=True)
    generated = 0
    skipped = 0
    failed: list[str] = []
    for index, student in enumerate(students, start=1):
        name = str(student.get("student_name") or "").strip()
        if not name:
            skipped += 1
            continue
        output = award_dir / f"{name}_奖状.png"
        if output.is_file():
            skipped += 1
            continue
        try:
            render_award(template, output, name)
            generated += 1
            print(f"  [{index}/{len(students)}] {name} -> {output.name}", flush=True)
        except Exception as error:
            failed.append(f"{name}: {error}")
            print(f"  [{index}/{len(students)}] {name} 失败：{error}", flush=True)

    print(f"完成：生成 {generated} 张，跳过 {skipped} 张，失败 {len(failed)} 张", flush=True)
    if failed:
        for item in failed[:20]:
            print(f"  FAIL {item}", flush=True)
    return 1 if failed and generated == 0 else 0


if __name__ == "__main__":
    try:
        raise SystemExit(main())
    except Exception as error:
        print(f"ERROR: {error}", file=sys.stderr)
        raise
