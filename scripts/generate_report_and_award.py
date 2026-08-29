#!/usr/bin/env python3
"""Generate wrong-question reports and awards from the workbench manifest.

Uses the workbench preview manifest as the single source of truth:
  - display score (protective score already applied when enabled);
  - recalculated wrong-question list (matches the display score);
  - band.

Outputs (same locations as the desktop tool):
  <source_dir>/全班错题报告/{姓名}_错题解析.pdf
  <source_dir>/已生成奖状/{姓名}_奖状.png   (only for score >= award_threshold)

Wrong-question reports embed the per-question explanation images
(<source_dir>/错题讲解/第N题.png). Existing files are regenerated only when
--force is given; otherwise they are skipped.
"""
import argparse
import json
import os
import sys
import tempfile
from pathlib import Path

try:
    sys.stdout.reconfigure(encoding="utf-8", errors="replace")
except Exception:
    pass

def resolve_cjk_font(source_dir: Path) -> Path | None:
    for candidate in (
        source_dir / "simhei.ttf",
        source_dir / "SimHei.ttf",
        Path(r"C:\Windows\Fonts\simhei.ttf"),
        Path(r"C:\Windows\Fonts\simsun.ttc"),
    ):
        if candidate.is_file():
            return candidate
    return None


def format_score(score):
    try:
        number = float(score)
        return str(int(number)) if number.is_integer() else str(number)
    except (TypeError, ValueError):
        return str(score or "")


def score_is_boundary(score) -> bool:
    """0 分和 100 分不生成错题报告；兼容整数/浮点/字符串成绩。"""
    try:
        number = float(score)
    except (TypeError, ValueError):
        return False
    return number in (0.0, 100.0)


def should_generate_report(student: dict) -> bool:
    """Skip reports for original/display 0 or 100, even after protective display scores."""
    if score_is_boundary(student.get("score")):
        return False
    if "original_score" in student and score_is_boundary(student.get("original_score")):
        return False
    return True


def generate_report(pdf_path: Path, student: dict, source_dir: Path) -> None:
    from fpdf import FPDF

    name = str(student.get("student_name") or "").strip()
    wrong_questions = [int(q) for q in (student.get("wrong_questions") or [])]
    wrong_count = int(student.get("wrong_count") or len(wrong_questions))

    pdf = FPDF(format="A4")
    pdf.set_auto_page_break(auto=True, margin=15)
    pdf.add_page()

    font_path = resolve_cjk_font(source_dir)
    if font_path is None:
        raise RuntimeError("找不到中文字体，请把 simhei.ttf 放到月考反馈文件夹根目录")
    pdf.add_font("cjk", "", str(font_path))
    pdf.set_font("cjk", "", 20)
    pdf.cell(0, 14, f"【{name}】专属错题解析报告", new_x="LMARGIN", new_y="NEXT", align="C")
    pdf.ln(4)

    pdf.ln(4)

    for question in wrong_questions:
        image = source_dir / "错题讲解" / f"第{question}题.png"
        pdf.set_font("cjk", "", 14)
        pdf.cell(0, 9, f"第 {question} 题", new_x="LMARGIN", new_y="NEXT")
        if image.is_file():
            pdf.image(str(image), w=170)
            pdf.ln(2)
        else:
            pdf.cell(0, 9, "（暂无该题讲解图）", new_x="LMARGIN", new_y="NEXT")
        pdf.ln(3)

    pdf.output(str(pdf_path))


def find_student_textbox(slide):
    def walk(container):
        for shape in container:
            if shape.shape_type == 6:
                result = walk(shape.shapes)
                if result is not None:
                    return result
            elif shape.has_text_frame:
                text = "".join(run.text for para in shape.text_frame.paragraphs for run in para.runs)
                if "学生" in text:
                    return shape
        return None
    return walk(slide.shapes)


def _wps_process_ids() -> set[str]:
    """Return running wpp.exe (WPS Presentation) PIDs, empty if wps not used."""
    try:
        import re
        import subprocess

        result = subprocess.run(
            ["tasklist", "/FI", "IMAGENAME eq wpp.exe", "/FO", "CSV", "/NH"],
            capture_output=True, text=True, timeout=10,
        )
        return {m.group(1) for m in re.finditer(r'"wpp\.exe","(\d+)"', result.stdout)}
    except Exception:
        return set()


def generate_award(award_path: Path, template_path: Path, student_name: str) -> None:
    from pptx import Presentation
    from pptx.util import Pt

    prs = Presentation(str(template_path))
    slide = prs.slides[0]
    box = find_student_textbox(slide)
    if box is None:
        raise RuntimeError("奖状模板中找不到学生姓名文本框")
    for para in box.text_frame.paragraphs:
        for run in para.runs:
            if "学生" in run.text:
                # 去除"学生："前缀，只保留学生姓名
                run.text = student_name
                # 名字字号自适应：模板原始字号约 190pt，替换后名字越长越减小，保证醒目且不溢出
                name_len = len(student_name)
                font_size = max(90, min(170, 170 - (name_len - 1) * 15))
                run.font.size = Pt(font_size)

    # 临时 pptx 放在系统临时目录（绝不写入奖状输出目录），导出后立即删除
    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False, prefix="award_tmp_") as tmp:
        tmp_path = Path(tmp.name)
    wps_before = _wps_process_ids()
    try:
        prs.save(str(tmp_path))
        import win32com.client

        app = win32com.client.DispatchEx("PowerPoint.Application")
        try:
            pres = app.Presentations.Open(str(tmp_path), WithWindow=False)
            try:
                pres.Slides(1).Export(str(award_path), "PNG", 1600, 1200)
            finally:
                pres.Close()
        finally:
            try:
                app.Quit()
            except Exception:
                pass
    finally:
        try:
            tmp_path.unlink(missing_ok=True)
        except OSError:
            pass

    # WPS 的 PowerPoint 兼容 COM 实例（wpp.exe）在 Quit 后可能残留，清理本次新启动的
    try:
        import subprocess

        leftover = _wps_process_ids() - wps_before
        for pid in leftover:
            try:
                subprocess.run(["taskkill", "/PID", str(pid), "/F"],
                               capture_output=True, timeout=10, check=False)
            except Exception:
                pass
    except Exception:
        pass


def main() -> int:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("--source-dir", required=True, type=Path)
    parser.add_argument("--manifest", required=True, type=Path)
    parser.add_argument("--award-threshold", type=float, default=80.0)
    parser.add_argument("--force", action="store_true", help="覆盖已存在的报告/奖状")
    parser.add_argument("--skip-reports", action="store_true")
    parser.add_argument("--skip-awards", action="store_true")
    parser.add_argument("--limit", type=int, default=0)
    args = parser.parse_args()

    source = args.source_dir.resolve()
    report_dir = source / "全班错题报告"
    award_dir = source / "已生成奖状"
    template = source / "优秀学员奖状1.pptx"
    report_dir.mkdir(parents=True, exist_ok=True)
    award_dir.mkdir(parents=True, exist_ok=True)

    # 清理历史残留的临时 pptx（桌面程序或异常中断可能在奖状目录留下 temp_*.pptx，且可能被 WPS 进程锁定）
    removed_temp = 0
    for stale in award_dir.glob("temp_*.pptx"):
        try:
            stale.unlink()
            removed_temp += 1
        except OSError:
            print(f"  警告：临时文件被占用，未能删除：{stale.name}", flush=True)
    if removed_temp:
        print(f"已清理 {removed_temp} 个历史临时文件（temp_*.pptx）", flush=True)

    manifest = json.loads(args.manifest.read_text(encoding="utf-8"))
    students = manifest.get("students") or []
    if args.limit > 0:
        students = students[: args.limit]

    reports = 0
    reports_skip = 0
    awards = 0
    awards_skip = 0
    failed: list[str] = []

    for index, student in enumerate(students, start=1):
        name = str(student.get("student_name") or "").strip()
        if not name:
            continue
        score = student.get("score")
        wrong_questions = [int(q) for q in (student.get("wrong_questions") or [])]

        # 错题报告：展示分 + 重算错题（有保护分就以保护分为准）；
        # 原始分或展示分为 0 / 100 的学生不生成错题报告。
        if not args.skip_reports:
            report_path = report_dir / f"{name}_错题解析.pdf"
            if not should_generate_report(student):
                reports_skip += 1
                print(f"  [{index}/{len(students)}] 跳过报告 {name}（{format_score(score)}分，0/100分不生成）", flush=True)
            elif report_path.is_file() and not args.force:
                reports_skip += 1
            else:
                try:
                    generate_report(report_path, student, source)
                    reports += 1
                except Exception as error:
                    failed.append(f"报告 {name}: {error}")
                print(f"  [{index}/{len(students)}] 报告 {name}（{format_score(score)}分/{len(wrong_questions)}题）", flush=True)

        # 奖状：展示分 >= 阈值
        if not args.skip_awards and score is not None and float(score) >= args.award_threshold:
            award_path = award_dir / f"{name}_奖状.png"
            if award_path.is_file() and not args.force:
                awards_skip += 1
            else:
                if not template.is_file():
                    failed.append(f"奖状 {name}: 模板缺失 {template}")
                else:
                    try:
                        generate_award(award_path, template, name)
                        awards += 1
                    except Exception as error:
                        failed.append(f"奖状 {name}: {error}")
            print(f"  [{index}/{len(students)}] 奖状 {name}", flush=True)

    print(
        f"完成：报告 生成 {reports} / 跳过 {reports_skip}；奖状 生成 {awards} / 跳过 {awards_skip}；失败 {len(failed)}",
        flush=True,
    )
    for item in failed[:20]:
        print(f"  FAIL {item}", flush=True)
    return 1 if failed and reports == 0 and awards == 0 else 0


if __name__ == "__main__":
    try:
        raise SystemExit(main())
    except Exception as error:
        print(f"ERROR: {error}", file=sys.stderr)
        raise
